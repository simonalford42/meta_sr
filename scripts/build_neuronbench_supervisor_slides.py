#!/usr/bin/env python3
"""Build the NeuronBench/evolve_pysr supervisor update as PPTX and PDF.

The deck is deliberately data-driven: the newest result summary is read from
run 708907, while the historical transfer rates come from the audited project
report.  No experiment is launched by this script.
"""

from __future__ import annotations

import json
import math
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import matplotlib

matplotlib.use("Agg")
import matplotlib.image as mpimg
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import Circle, FancyBboxPatch, Polygon, Rectangle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "presentations"
STEM = "neuronbench_evolve_pysr_supervisor_2026-08-28"
ASSET_DIR = OUT_DIR / f"{STEM}_assets"
PPTX_PATH = OUT_DIR / f"{STEM}.pptx"
PDF_PATH = OUT_DIR / f"{STEM}.pdf"
HH_IMAGE = Path("/home/sca63/sandia_spiking/Hodgkins-Huxley/hh_response2.png")

W, H = 13.333, 7.5

# Sandia-inspired technical palette: restrained navy with teal/orange evidence accents.
NAVY = "17365D"
INK = "17212B"
MUTED = "5F6B76"
BG = "F7F9FB"
WHITE = "FFFFFF"
SOFT = "EAF0F5"
TEAL = "168C8C"
TEAL_DARK = "0D666A"
TEAL_LIGHT = "DDF2F0"
ORANGE = "E7812B"
ORANGE_LIGHT = "FCE9D8"
BLUE = "4D79A7"
BLUE_LIGHT = "DDE9F4"
GREEN = "4F9D69"
GREEN_LIGHT = "E2F1E7"
RED = "C9584E"
RED_LIGHT = "F6E2E0"
GRAY = "D5DDE5"


def rgb(hex_color: str) -> tuple[float, float, float]:
    return tuple(int(hex_color[i : i + 2], 16) / 255 for i in (0, 2, 4))


def ppt_rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


@dataclass
class Layer:
    kind: str
    x: float
    y: float
    w: float = 0.0
    h: float = 0.0
    text: str = ""
    fill: str | None = None
    line: str | None = None
    lw: float = 1.0
    size: float = 18.0
    color: str = INK
    bold: bool = False
    italic: bool = False
    align: str = "left"
    valign: str = "top"
    radius: float = 0.10
    path: str | None = None
    alpha: float = 1.0
    points: list[tuple[float, float]] | None = None


@dataclass
class SlideSpec:
    title: str
    kicker: str = ""
    footer: str = ""
    dark: bool = False
    layers: list[Layer] = field(default_factory=list)

    def add(self, kind: str, x: float, y: float, w: float = 0, h: float = 0, **kwargs: Any) -> Layer:
        layer = Layer(kind=kind, x=x, y=y, w=w, h=h, **kwargs)
        self.layers.append(layer)
        return layer


def add_standard_header(slide: SlideSpec, number: int) -> None:
    if slide.dark:
        slide.add("rect", 0, 0, W, H, fill=NAVY, line=NAVY)
        slide.add("rect", 0, 0, 0.13, H, fill=TEAL, line=TEAL)
        if slide.kicker:
            slide.add("text", 0.70, 0.52, 11.8, 0.28, text=slide.kicker.upper(), size=10, color="9ED9D6", bold=True)
        slide.add("text", 0.70, 0.92, 11.8, 1.20, text=slide.title, size=31, color=WHITE, bold=True, valign="center")
        slide.add("text", 12.25, 0.46, 0.45, 0.30, text=f"{number:02d}", size=10, color="9ED9D6", bold=True, align="right")
    else:
        slide.add("rect", 0, 0, W, H, fill=BG, line=BG)
        slide.add("rect", 0, 0, W, 0.12, fill=TEAL, line=TEAL)
        if slide.kicker:
            slide.add("text", 0.66, 0.43, 11.9, 0.23, text=slide.kicker.upper(), size=9.5, color=TEAL_DARK, bold=True)
        slide.add("text", 0.66, 0.75, 11.9, 0.58, text=slide.title, size=26, color=NAVY, bold=True)
        slide.add("text", 12.24, 0.43, 0.45, 0.25, text=f"{number:02d}", size=9.5, color=MUTED, bold=True, align="right")
    if slide.footer:
        footer_color = "B8C8D8" if slide.dark else MUTED
        slide.add("text", 0.66, 7.14, 12.0, 0.20, text=slide.footer, size=7.2, color=footer_color)


def card(slide: SlideSpec, x: float, y: float, w: float, h: float, *, fill: str = WHITE, line: str = GRAY, radius: float = 0.12) -> None:
    slide.add("roundrect", x, y, w, h, fill=fill, line=line, lw=0.8, radius=radius)


def pill(slide: SlideSpec, x: float, y: float, w: float, text: str, *, fill: str, color: str = WHITE) -> None:
    slide.add("roundrect", x, y, w, 0.38, fill=fill, line=fill, radius=0.18)
    slide.add("text", x, y + 0.03, w, 0.28, text=text, size=10.5, color=color, bold=True, align="center", valign="center")


def arrow(slide: SlideSpec, x1: float, y1: float, x2: float, y2: float, *, color: str = TEAL, lw: float = 2.0) -> None:
    slide.add("line", x1, y1, x2 - x1, y2 - y1, line=color, lw=lw)
    angle = math.atan2(y2 - y1, x2 - x1)
    d = 0.16
    wing = 0.075
    tip = (x2, y2)
    base = (x2 - d * math.cos(angle), y2 - d * math.sin(angle))
    p1 = (base[0] + wing * math.sin(angle), base[1] - wing * math.cos(angle))
    p2 = (base[0] - wing * math.sin(angle), base[1] + wing * math.cos(angle))
    slide.add("poly", 0, 0, points=[tip, p1, p2], fill=color, line=color)


def bullet(slide: SlideSpec, x: float, y: float, w: float, title: str, body: str, *, color: str = TEAL, body_size: float = 15.5) -> None:
    slide.add("circle", x, y + 0.10, 0.13, 0.13, fill=color, line=color)
    slide.add("text", x + 0.27, y, w - 0.27, 0.30, text=title, size=body_size, color=INK, bold=True)
    slide.add("text", x + 0.27, y + 0.34, w - 0.27, 0.55, text=body, size=body_size - 2.0, color=MUTED)


def build_slides(new_summary: dict[str, Any]) -> list[SlideSpec]:
    slides: list[SlideSpec] = []

    # 1 — Title / executive signal
    s = SlideSpec(
        title="NeuronBench × evolve_pysr",
        kicker="Grant update • 28 August 2026",
        footer="Project sources: meta_sr runs 313196, 313195, 190177, 708907 • NeuronBench arXiv:2608.09696v4",
        dark=True,
    )
    add_standard_header(s, 1)
    s.add("text", 0.72, 2.10, 11.3, 0.55, text="Meta-evolving symbolic regression for mechanistic neuron dynamics", size=19, color="D7E2EC")
    stats = [
        ("6", "mechanistically distinct\nneuron worlds", TEAL),
        ("4", "PySR search components\nco-evolved", ORANGE),
        ("23/25", "strict held-out recoveries\nin the newest run", GREEN),
    ]
    for i, (big, label, accent) in enumerate(stats):
        x = 0.74 + i * 4.05
        s.add("roundrect", x, 3.14, 3.62, 1.58, fill="21486F", line="416886", lw=0.8, radius=0.14)
        s.add("rect", x, 3.14, 0.09, 1.58, fill=accent, line=accent)
        s.add("text", x + 0.32, 3.39, 2.95, 0.55, text=big, size=31, color=WHITE, bold=True)
        s.add("text", x + 0.32, 4.05, 3.02, 0.48, text=label, size=12.5, color="D7E2EC")
    pill(s, 0.75, 5.27, 2.25, "NEW: DOMAIN-BLIND", fill=TEAL)
    s.add("text", 3.23, 5.25, 8.75, 0.70, text="A generic prompt still discovered the same useful search motif:\nfit the residual, then add the most informative feature direction.", size=17.5, color=WHITE, bold=True)
    slides.append(s)

    # 2 — Original NeuronBench
    s = SlideSpec(
        title="NeuronBench asks for active mechanistic discovery",
        kicker="Kevin Murphy’s Model Discovery Agent paper",
        footer="Source: Murphy, Model Discovery Agent, arXiv:2608.09696v4 (25 Aug 2026), §4.4 and Appendix E",
    )
    add_standard_header(s, 2)
    pill(s, 0.68, 1.48, 2.15, "ORIGINAL TASK", fill=NAVY)
    s.add("text", 3.06, 1.43, 9.45, 0.53, text="Six “mystery neurons” built from generalized Hodgkin–Huxley ODEs", size=18, color=INK, bold=True)
    # Pipeline
    boxes = [
        (0.72, 2.30, 2.55, "Choose an experiment", "9 current-clamp templates\n(active design)"),
        (3.91, 2.30, 2.55, "Probe a mystery neuron", "novel membrane mechanism\n+ hidden gating state"),
        (7.10, 2.30, 2.55, "Observe a trajectory", "voltage over time\n(+ stochastic variant)"),
        (10.29, 2.30, 2.32, "Predict interventions", "6 trajectory features\non held-out probes"),
    ]
    for i, (x, y, w, head, body) in enumerate(boxes):
        card(s, x, y, w, 1.52, fill=WHITE)
        s.add("circle", x + 0.18, y + 0.20, 0.40, 0.40, fill=[TEAL, ORANGE, BLUE, GREEN][i], line=[TEAL, ORANGE, BLUE, GREEN][i])
        s.add("text", x + 0.18, y + 0.24, 0.40, 0.24, text=str(i + 1), size=11, color=WHITE, bold=True, align="center")
        s.add("text", x + 0.72, y + 0.18, w - 0.88, 0.42, text=head, size=14.5, color=NAVY, bold=True)
        s.add("text", x + 0.23, y + 0.80, w - 0.44, 0.52, text=body, size=12.2, color=MUTED, align="center")
        if i < 3:
            arrow(s, x + w + 0.15, y + 0.76, boxes[i + 1][0] - 0.15, y + 0.76, color=TEAL_DARK, lw=1.8)
    s.add("text", 0.78, 4.40, 3.25, 0.34, text="Why it is hard", size=16, color=NAVY, bold=True)
    difficulties = [
        (0.78, "Partial observability", "Infer hidden gates and mechanisms\nfrom voltage alone.", ORANGE),
        (4.58, "Experiment choice", "Select useful perturbations\nunder a small budget.", TEAL),
        (8.38, "Open hypothesis space", "The novel mechanism may be\nabsent from the model pool.", BLUE),
    ]
    for x, head, body, accent in difficulties:
        card(s, x, 4.83, 3.55, 0.98, fill=WHITE)
        s.add("circle", x + 0.19, 5.05, 0.18, 0.18, fill=accent, line=accent)
        s.add("text", x + 0.50, 4.97, 2.80, 0.28, text=head, size=13.2, color=INK, bold=True)
        s.add("text", x + 0.50, 5.30, 2.80, 0.38, text=body, size=10.5, color=MUTED)
    s.add("roundrect", 0.76, 6.18, 11.90, 0.58, fill=SOFT, line=SOFT, radius=0.10)
    s.add("text", 1.02, 6.31, 11.40, 0.27, text="MDA combines nested SMC³ inference + LLM model proposals + Value-of-Information experiment design.", size=14.2, color=NAVY, bold=True, align="center")
    slides.append(s)

    # 3 — Reduction
    s = SlideSpec(
        title="We isolate the equation-discovery problem",
        kicker="Our adaptation: NeuronBench → ordinary symbolic regression",
        footer="Implementation: scripts/neuronbench_fully_observable.py • domain: domains.NeuronBenchDomain",
    )
    add_standard_header(s, 3)
    s.add("text", 0.72, 1.40, 5.35, 0.40, text="Original benchmark", size=15.5, color=MUTED, bold=True)
    s.add("text", 7.20, 1.40, 5.35, 0.40, text="Fully observable SR benchmark", size=15.5, color=TEAL_DARK, bold=True)
    left = ["choose current protocol", "observe voltage trajectory", "infer hidden gates + mechanism", "forecast held-out interventions"]
    right = ["sample independent neuron states", "observe I_ext, V, and each channel-open fraction φ_c", "fit one scalar target: dV/dt", "evaluate exact physical vector field on held-out states"]
    for i, txt in enumerate(left):
        card(s, 0.74, 1.90 + 0.84 * i, 4.92, 0.60, fill=WHITE)
        s.add("text", 1.02, 2.04 + 0.84 * i, 4.37, 0.26, text=txt, size=13.5, color=INK)
    for i, txt in enumerate(right):
        card(s, 7.18, 1.90 + 0.84 * i, 5.40, 0.60, fill=TEAL_LIGHT, line="A7D6D2")
        s.add("text", 7.47, 2.04 + 0.84 * i, 4.83, 0.28, text=txt, size=13.2, color=TEAL_DARK, bold=(i == 2))
    arrow(s, 5.86, 3.15, 6.89, 3.15, color=ORANGE, lw=2.8)
    pill(s, 5.93, 2.54, 0.88, "REDUCE", fill=ORANGE)
    s.add("roundrect", 1.25, 5.52, 10.83, 0.92, fill=NAVY, line=NAVY, radius=0.15)
    s.add("text", 1.55, 5.71, 10.20, 0.45, text="dV/dt = I_ext − Σ_c g_c φ_c (V − E_c)", size=22, color=WHITE, bold=True, align="center")
    fact_data = [("1,024", "train states"), ("16,384", "held-out states"), ("0", "observation noise"), ("{+, −, ×}", "operator set")]
    for i, (n, lab) in enumerate(fact_data):
        x = 1.35 + i * 2.73
        s.add("text", x, 6.58, 1.95, 0.28, text=n, size=16.5, color=NAVY, bold=True, align="center")
        s.add("text", x, 6.88, 1.95, 0.18, text=lab, size=9.5, color=MUTED, align="center")
    slides.append(s)

    # 4 — HH bridge
    s = SlideSpec(
        title="Same Hodgkin–Huxley skeleton, broader discovery test",
        kicker="Connection to ~/sandia_spiking",
        footer="Existing HH artifact: /home/sca63/sandia_spiking/Hodgkins-Huxley/hh.pptx and hh_response2.png",
    )
    add_standard_header(s, 4)
    card(s, 0.68, 1.48, 5.20, 4.98, fill=WHITE)
    s.add("image", 0.91, 1.74, 4.74, 4.21, path=str(HH_IMAGE))
    s.add("text", 1.05, 6.02, 4.46, 0.24, text="Existing project: identify HH dynamics from simulated trajectories", size=11.5, color=MUTED, align="center")
    s.add("text", 6.30, 1.54, 6.20, 0.45, text="The shared physical law", size=17, color=NAVY, bold=True)
    s.add("roundrect", 6.29, 2.08, 6.10, 1.05, fill=NAVY, line=NAVY, radius=0.12)
    s.add("text", 6.55, 2.33, 5.58, 0.48, text="C_m dV/dt = I_ext − I_Na − I_K − I_leak − I_novel", size=18, color=WHITE, bold=True, align="center")
    s.add("line", 9.33, 3.40, 0, 2.36, line=GRAY, lw=1.2)
    card(s, 6.26, 3.61, 2.85, 1.76, fill=BLUE_LIGHT, line="ABC3DB")
    s.add("text", 6.53, 3.84, 2.30, 0.31, text="sandia_spiking", size=15.2, color=NAVY, bold=True, align="center")
    s.add("text", 6.53, 4.30, 2.30, 0.76, text="one canonical HH system\ntrajectory data\nrecover coupled ODEs", size=12.5, color=MUTED, align="center")
    card(s, 9.58, 3.61, 2.85, 1.76, fill=TEAL_LIGHT, line="A7D6D2")
    s.add("text", 9.83, 3.84, 2.35, 0.31, text="NeuronBench SR", size=15.2, color=TEAL_DARK, bold=True, align="center")
    s.add("text", 9.83, 4.30, 2.35, 0.76, text="six mechanisms\ncollocation states\nrecover membrane law", size=12.5, color=MUTED, align="center")
    arrow(s, 9.17, 4.49, 9.49, 4.49, color=ORANGE, lw=2.2)
    s.add("roundrect", 6.28, 5.81, 6.12, 0.72, fill=ORANGE_LIGHT, line="F0C49C", radius=0.10)
    s.add("text", 6.57, 5.98, 5.55, 0.31, text="NeuronBench turns the HH demonstration into a transfer benchmark.", size=14.5, color="8B4B17", bold=True, align="center")
    slides.append(s)

    # 5 — Six problems
    s = SlideSpec(
        title="The benchmark is six current-balance problems",
        kicker="Common Na + K + leak backbone; one world-specific mechanism",
        footer="World definitions: neuronbench.worlds at commit c354622458c460b419cab821d482c879f0578377",
    )
    add_standard_header(s, 5)
    problems = [
        ("Z-rebound", "φ_Z = m_Z²h_Z", "−4 φ_Z (V − 120)", "depolarization block", TEAL),
        ("H-sag", "φ_h = m_h", "−5 φ_h (V + 30)", "sag + rebound", BLUE),
        ("Na-fatigue", "φ_Na includes s_Na", "no extra current term", "use-dependent rundown", ORANGE),
        ("Ca-rebound", "φ_T = m_T²h_T", "−3.2 φ_T (V − 120)", "rebound burst", RED),
        ("D-type", "φ_D = m_Dh_D", "−9 φ_D (V + 77)", "delayed firing", GREEN),
        ("Textbook-M", "φ_M = m_M", "−2.5 φ_M (V + 77)", "spike adaptation", "8064A2"),
    ]
    for i, (name, state, term, phenotype, accent) in enumerate(problems):
        col = i % 3
        row = i // 3
        x = 0.70 + col * 4.17
        y = 1.53 + row * 2.45
        card(s, x, y, 3.78, 2.09, fill=WHITE)
        s.add("rect", x, y, 0.10, 2.09, fill=accent, line=accent)
        s.add("text", x + 0.31, y + 0.20, 3.14, 0.35, text=name, size=17, color=NAVY, bold=True)
        pill(s, x + 2.60, y + 0.17, 0.84, f"{i + 1}/6", fill=accent)
        s.add("text", x + 0.31, y + 0.76, 3.10, 0.28, text=state, size=13, color=INK, bold=True)
        s.add("text", x + 0.31, y + 1.14, 3.10, 0.26, text=term, size=12.5, color=TEAL_DARK)
        s.add("text", x + 0.31, y + 1.60, 3.10, 0.24, text=phenotype, size=11, color=MUTED, italic=True)
    s.add("roundrect", 0.72, 6.48, 12.00, 0.43, fill=SOFT, line=SOFT, radius=0.08)
    s.add("text", 0.96, 6.57, 11.52, 0.21, text="Ground-truth complexity: 23 nodes for Na-fatigue; 31 nodes for the five worlds with an extra current.", size=11.5, color=NAVY, align="center")
    slides.append(s)

    # 6 — evolve_pysr
    s = SlideSpec(
        title="evolve_pysr searches over the search algorithm",
        kicker="Meta-evolution, not just equation evolution",
        footer="Evolution fitness: fraction of training worlds solved at held-out NRMSE ≤ 1e−6 • final fits use up to 1e6 evaluations",
    )
    add_standard_header(s, 6)
    s.add("text", 0.72, 1.44, 3.25, 0.35, text="LLM proposes Julia operators", size=16, color=NAVY, bold=True)
    ops = [("Mutation", TEAL), ("Survival", ORANGE), ("Selection", BLUE), ("Loss", GREEN)]
    for i, (name, accent) in enumerate(ops):
        y = 1.98 + i * 0.79
        card(s, 0.72, y, 2.75, 0.58, fill=WHITE)
        s.add("rect", 0.72, y, 0.10, 0.58, fill=accent, line=accent)
        s.add("text", 1.02, y + 0.15, 2.15, 0.26, text=name, size=14.2, color=INK, bold=True)
    arrow(s, 3.66, 3.31, 5.06, 3.31, color=TEAL_DARK, lw=2.4)
    card(s, 5.15, 2.16, 3.10, 2.32, fill=NAVY, line=NAVY)
    s.add("text", 5.50, 2.45, 2.40, 0.36, text="PySR evaluation", size=18, color=WHITE, bold=True, align="center")
    s.add("text", 5.53, 3.05, 2.34, 0.90, text="3 seeds / candidate\n10 seeds / finalists\nscore = solve fraction", size=13, color="D7E2EC", align="center")
    arrow(s, 8.43, 3.31, 9.78, 3.31, color=TEAL_DARK, lw=2.4)
    card(s, 9.88, 2.16, 2.74, 2.32, fill=TEAL_LIGHT, line="A7D6D2")
    s.add("text", 10.18, 2.48, 2.16, 0.36, text="Select top-k", size=18, color=TEAL_DARK, bold=True, align="center")
    s.add("text", 10.18, 3.07, 2.16, 0.78, text="retain strong bundles\nmutate + recombine\nrepeat by generation", size=12.7, color=MUTED, align="center")
    arrow(s, 11.25, 4.78, 2.10, 4.78, color=ORANGE, lw=1.7)
    s.add("text", 0.72, 5.22, 4.20, 0.31, text="Transfer regimes", size=15.5, color=NAVY, bold=True)
    regimes = [("top-1", "train on Z", TEAL), ("top-2", "train on Z + H", BLUE), ("“top5”", "train on five; hold out one", ORANGE)]
    for i, (name, desc, accent) in enumerate(regimes):
        x = 0.72 + i * 2.53
        card(s, x, 5.68, 2.25, 0.92, fill=WHITE)
        s.add("text", x + 0.18, 5.84, 1.88, 0.25, text=name, size=14, color=accent, bold=True, align="center")
        s.add("text", x + 0.18, 6.17, 1.88, 0.20, text=desc, size=9.7, color=MUTED, align="center")
    s.add("roundrect", 8.67, 5.35, 3.95, 1.24, fill=ORANGE_LIGHT, line="F0C49C", radius=0.12)
    s.add("text", 8.94, 5.56, 3.40, 0.30, text="Repeated evolved motif", size=13, color="8B4B17", bold=True, align="center")
    s.add("text", 8.94, 5.98, 3.40, 0.37, text="project residual onto a feature\nthen add the fitted correction", size=12.5, color=INK, align="center")
    slides.append(s)

    # 7 — historical transfer
    s = SlideSpec(
        title="Earlier evolution transferred across unseen neurons",
        kicker="Whole-frontier manual structural recovery",
        footer="Audited sources: reports/neuron_topk_manual_transfer_report.pdf and neuron_manual_match_comparison.json",
    )
    add_standard_header(s, 7)
    s.add("text", 0.72, 1.43, 7.10, 0.38, text="Manual match rate on held-out worlds", size=16, color=NAVY, bold=True)
    rates = [("top-1", 19, 25, 0.76, TEAL), ("top-2", 16, 20, 0.80, BLUE), ("“top5”", 5, 5, 1.00, ORANGE)]
    for i, (name, matches, total, rate, accent) in enumerate(rates):
        y = 2.02 + i * 1.16
        s.add("text", 0.78, y + 0.10, 1.13, 0.27, text=name, size=14, color=INK, bold=True)
        s.add("roundrect", 2.00, y, 5.48, 0.55, fill=SOFT, line=SOFT, radius=0.10)
        s.add("roundrect", 2.00, y, 5.48 * rate, 0.55, fill=accent, line=accent, radius=0.10)
        s.add("text", 2.19, y + 0.12, 4.95, 0.24, text=f"{matches}/{total}", size=12.5, color=WHITE, bold=True)
        s.add("text", 7.70, y + 0.07, 0.78, 0.31, text=f"{rate:.0%}", size=17, color=accent, bold=True, align="right")
    card(s, 8.78, 1.62, 3.82, 3.99, fill=WHITE)
    s.add("text", 9.08, 1.92, 3.20, 0.33, text="What counted as recovery?", size=16, color=NAVY, bold=True, align="center")
    checks = [
        "all required physical monomials",
        "no material extra monomial",
        "clearly close coefficients",
        "tiny numerical artifacts allowed",
    ]
    for i, txt in enumerate(checks):
        y = 2.58 + i * 0.62
        s.add("circle", 9.10, y + 0.04, 0.23, 0.23, fill=GREEN_LIGHT, line=GREEN)
        s.add("text", 9.11, y + 0.055, 0.21, 0.14, text=str(i + 1), size=8.5, color=GREEN, bold=True, align="center")
        s.add("text", 9.49, y, 2.73, 0.31, text=txt, size=12.2, color=INK)
    s.add("roundrect", 0.74, 5.66, 7.72, 0.76, fill=TEAL_LIGHT, line="A7D6D2", radius=0.10)
    s.add("text", 1.05, 5.86, 7.13, 0.34, text="Cross-run signal: residual-guided additive mutation appears in all three best bundles.", size=14, color=TEAL_DARK, bold=True, align="center")
    s.add("roundrect", 8.78, 5.82, 3.82, 0.60, fill=ORANGE_LIGHT, line="F0C49C", radius=0.10)
    s.add("text", 9.03, 5.95, 3.30, 0.28, text="Caveat: “top5” = one LOOCV fold", size=11.5, color="8B4B17", bold=True, align="center")
    slides.append(s)

    # 8 — newest uninformative
    s = SlideSpec(
        title="Newest result: domain-blind prompts still transfer",
        kicker="Run 708907 • “uninformative prompt” • train on Z-rebound only",
        footer="Source: runs/708907/run_data.json and runs/708907/neuron_full_eval/neuron_results.json • five final seeds/world",
    )
    add_standard_header(s, 8)
    s.add("roundrect", 0.71, 1.46, 5.08, 1.14, fill=NAVY, line=NAVY, radius=0.12)
    s.add("text", 0.99, 1.67, 4.51, 0.66, text="“Improve the algorithm’s ability to discover\nthe expression that generated the task data.”", size=15.2, color=WHITE, italic=True, align="center")
    s.add("text", 0.73, 2.85, 5.00, 0.30, text="Best training solve rate by generation", size=14.5, color=NAVY, bold=True)
    vals = [0.40, 2 / 3, 1.0] + [1.0] * 12
    x0, y0, cw, ch = 0.90, 5.12, 4.66, 1.75
    s.add("line", x0, y0, cw, 0, line=GRAY, lw=1)
    s.add("line", x0, y0 - ch, 0, ch, line=GRAY, lw=1)
    pts = []
    for i, val in enumerate(vals):
        x = x0 + cw * i / 14
        y = y0 - ch * val
        pts.append((x, y))
    for a, b in zip(pts[:-1], pts[1:]):
        s.add("line", a[0], a[1], b[0] - a[0], b[1] - a[1], line=TEAL, lw=2.4)
    for i in [0, 1, 2, 14]:
        x, y = pts[i]
        s.add("circle", x - 0.055, y - 0.055, 0.11, 0.11, fill=TEAL, line=WHITE, lw=0.6)
    for tick, val in [(0, "0"), (0.5, "50%"), (1, "100%")]:
        yy = y0 - ch * tick
        s.add("text", 0.42, yy - 0.09, 0.40, 0.18, text=val, size=8.5, color=MUTED, align="right")
    for gen, xx in [(1, x0), (3, pts[2][0]), (15, x0 + cw)]:
        s.add("text", xx - 0.22, y0 + 0.12, 0.44, 0.17, text=str(gen), size=8.5, color=MUTED, align="center")
    s.add("text", 2.65, 5.48, 1.20, 0.20, text="generation", size=9, color=MUTED, align="center")
    s.add("line", x0, y0, cw, 0, line=RED, lw=1.0)
    s.add("text", 1.02, 5.17, 1.24, 0.18, text="baseline = 0%", size=8.8, color=RED)

    # Held-out results
    s.add("text", 6.18, 1.47, 6.25, 0.36, text="Held-out strict numerical recovery", size=16, color=NAVY, bold=True)
    per_world = new_summary["per_world"]
    ordered = ["h_sag", "na_fatigue", "ca_rebound", "d_type", "textbook_M"]
    labels = {"h_sag": "H-sag", "na_fatigue": "Na-fatigue", "ca_rebound": "Ca-rebound", "d_type": "D-type", "textbook_M": "Textbook-M"}
    for i, world in enumerate(ordered):
        rec = per_world[world]["recovered"]
        near = per_world[world]["near"]
        y = 2.00 + i * 0.68
        s.add("text", 6.22, y + 0.07, 1.30, 0.24, text=labels[world], size=11.5, color=INK, bold=True)
        for j in range(5):
            fill = GREEN if j < rec else BLUE
            s.add("roundrect", 7.72 + j * 0.55, y, 0.42, 0.42, fill=fill, line=WHITE, lw=0.8, radius=0.06)
        med = per_world[world]["median"]
        s.add("text", 10.72, y + 0.07, 1.78, 0.24, text=f"median {med:.1e}", size=10, color=MUTED, align="right")
    s.add("text", 7.73, 5.48, 3.02, 0.19, text="■ strict recovered     ■ near-exact", size=9.5, color=MUTED)
    s.add("roundrect", 6.18, 5.82, 6.21, 0.83, fill=GREEN_LIGHT, line="B9DCC4", radius=0.12)
    s.add("text", 6.48, 5.98, 1.60, 0.38, text="23 / 25", size=24, color=GREEN, bold=True)
    s.add("text", 8.19, 5.98, 3.93, 0.40, text="strict recoveries\n+ 2 near-exact, 0 misses", size=13, color=INK, bold=True)
    pill(s, 0.74, 6.30, 2.05, "BEST BUNDLE", fill=ORANGE)
    s.add("text", 3.02, 6.19, 2.79, 0.55, text="residual projection mutation\n+ frequency-aware selection + MSE", size=11.5, color=INK, bold=True)
    slides.append(s)

    # 9 — Takeaways
    s = SlideSpec(
        title="The evidence now supports a reusable search principle",
        kicker="Takeaway for the program",
        footer="Recommended next evidence: replicate seeds/runs, complete all six top5 LOOCV folds, and manually audit run 708907 equations",
        dark=True,
    )
    add_standard_header(s, 9)
    takeaways = [
        ("01", "Mechanistic benchmark", "We preserved the exact neuron current balance while turning active, latent discovery\ninto a controlled symbolic-regression test."),
        ("02", "Transfer, not memorization", "Algorithms evolved on 1–5 neuron worlds recover equations on unseen mechanisms;\nthe newest run used a generic prompt."),
        ("03", "A stable algorithmic motif", "Residual-guided additive mutation repeatedly emerges because neuron currents\ncombine as additive feature directions."),
    ]
    for i, (num, head, body) in enumerate(takeaways):
        y = 2.15 + i * 1.20
        s.add("text", 0.78, y, 0.65, 0.40, text=num, size=15, color="9ED9D6", bold=True)
        s.add("text", 1.55, y - 0.03, 3.14, 0.35, text=head, size=17, color=WHITE, bold=True)
        s.add("text", 4.85, y - 0.02, 7.40, 0.70, text=body, size=12.8, color="D7E2EC")
        if i < 2:
            s.add("line", 0.78, y + 0.83, 11.48, 0, line="416886", lw=0.8)
    s.add("roundrect", 0.78, 5.93, 11.48, 0.73, fill=TEAL, line=TEAL, radius=0.14)
    s.add("text", 1.14, 6.12, 10.76, 0.34, text="Bottom line: meta-evolution is learning a useful search heuristic—not merely receiving neuron-specific hints.", size=16, color=WHITE, bold=True, align="center")
    slides.append(s)

    return slides


def render_matplotlib(slide: SlideSpec, page_num: int, pdf: PdfPages) -> None:
    fig = plt.figure(figsize=(W, H), dpi=150)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, W)
    ax.set_ylim(H, 0)
    ax.axis("off")
    fig.patch.set_facecolor(rgb(NAVY if slide.dark else BG))

    for l in slide.layers:
        edge = rgb(l.line) if l.line else "none"
        face = rgb(l.fill) if l.fill else "none"
        if l.kind == "rect":
            ax.add_patch(Rectangle((l.x, l.y), l.w, l.h, facecolor=face, edgecolor=edge, linewidth=l.lw, alpha=l.alpha))
        elif l.kind == "roundrect":
            ax.add_patch(FancyBboxPatch((l.x, l.y), l.w, l.h, boxstyle=f"round,pad=0.01,rounding_size={l.radius}", facecolor=face, edgecolor=edge, linewidth=l.lw, alpha=l.alpha))
        elif l.kind == "circle":
            ax.add_patch(Circle((l.x + l.w / 2, l.y + l.h / 2), min(l.w, l.h) / 2, facecolor=face, edgecolor=edge, linewidth=l.lw, alpha=l.alpha))
        elif l.kind == "line":
            ax.plot([l.x, l.x + l.w], [l.y, l.y + l.h], color=rgb(l.line or INK), linewidth=l.lw, solid_capstyle="round")
        elif l.kind == "poly":
            ax.add_patch(Polygon(l.points or [], closed=True, facecolor=face, edgecolor=edge, linewidth=l.lw))
        elif l.kind == "image":
            img = mpimg.imread(l.path)
            ax.imshow(
                img[::-1],
                extent=(l.x, l.x + l.w, l.y, l.y + l.h),
                aspect="auto",
                origin="upper",
                zorder=2,
            )
        elif l.kind == "text":
            ha = {"left": "left", "center": "center", "right": "right"}[l.align]
            va = {"top": "top", "center": "center", "bottom": "bottom"}.get(l.valign, "top")
            tx = l.x if l.align == "left" else (l.x + l.w / 2 if l.align == "center" else l.x + l.w)
            ty = l.y if l.valign == "top" else (l.y + l.h / 2 if l.valign == "center" else l.y + l.h)
            ax.text(tx, ty, l.text, ha=ha, va=va, fontsize=l.size, color=rgb(l.color), weight="bold" if l.bold else "normal", style="italic" if l.italic else "normal", family="Liberation Sans", linespacing=1.18)

    # Image artists can re-enable autoscaling; restore the slide coordinate box.
    ax.set_xlim(0, W)
    ax.set_ylim(H, 0)
    png_path = ASSET_DIR / f"slide_{page_num:02d}.png"
    fig.savefig(png_path, dpi=150, facecolor=fig.get_facecolor())
    pdf.savefig(fig, facecolor=fig.get_facecolor())
    plt.close(fig)


def ppt_add_text(slide, l: Layer) -> None:
    shape = slide.shapes.add_textbox(Inches(l.x), Inches(l.y), Inches(max(l.w, 0.05)), Inches(max(l.h, 0.05)))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "center": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}.get(l.valign, MSO_ANCHOR.TOP)
    lines = l.text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[l.align]
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.font.name = "Arial"
        p.font.size = Pt(l.size)
        p.font.bold = l.bold
        p.font.italic = l.italic
        p.font.color.rgb = ppt_rgb(l.color)


def render_pptx(slides: list[SlideSpec]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    for spec in slides:
        slide = prs.slides.add_slide(blank)
        for l in spec.layers:
            if l.kind == "text":
                ppt_add_text(slide, l)
                continue
            if l.kind == "image":
                slide.shapes.add_picture(l.path, Inches(l.x), Inches(l.y), Inches(l.w), Inches(l.h))
                continue
            if l.kind == "line":
                shp = slide.shapes.add_connector(1, Inches(l.x), Inches(l.y), Inches(l.x + l.w), Inches(l.y + l.h))
                shp.line.color.rgb = ppt_rgb(l.line or INK)
                shp.line.width = Pt(l.lw)
                continue
            if l.kind == "poly":
                pts = l.points or []
                if len(pts) == 3:
                    xs = [p[0] for p in pts]
                    ys = [p[1] for p in pts]
                    shp = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(min(xs)), Inches(min(ys)), Inches(max(xs) - min(xs)), Inches(max(ys) - min(ys)))
                    tip = pts[0]
                    base = ((pts[1][0] + pts[2][0]) / 2, (pts[1][1] + pts[2][1]) / 2)
                    angle = math.degrees(math.atan2(tip[1] - base[1], tip[0] - base[0]))
                    shp.rotation = 90 + angle
                else:
                    continue
            else:
                shape_type = {
                    "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    "roundrect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                    "circle": MSO_AUTO_SHAPE_TYPE.OVAL,
                }[l.kind]
                shp = slide.shapes.add_shape(shape_type, Inches(l.x), Inches(l.y), Inches(max(l.w, 0.01)), Inches(max(l.h, 0.01)))
            shp.fill.solid()
            shp.fill.fore_color.rgb = ppt_rgb(l.fill or WHITE)
            shp.fill.transparency = int((1 - l.alpha) * 100)
            if l.line:
                shp.line.color.rgb = ppt_rgb(l.line)
                shp.line.width = Pt(l.lw)
            else:
                shp.line.fill.background()
    prs.save(PPTX_PATH)


def load_newest_summary() -> dict[str, Any]:
    path = ROOT / "runs/708907/neuron_full_eval/neuron_results.json"
    with path.open(encoding="utf-8") as handle:
        payload = json.load(handle)
    per_world: dict[str, Any] = {}
    for world, info in payload["per_world"].items():
        per_world[world] = {
            "recovered": int(info["counts"]["recovered"]),
            "near": int(info["counts"]["near-exact"]),
            "median": float(info["median_best_nrmse"]),
        }
    return {"per_world": per_world}


def write_readme() -> None:
    readme = OUT_DIR / f"{STEM}_README.md"
    text = f"""# NeuronBench × evolve_pysr supervisor slides

- PowerPoint: `{PPTX_PATH.name}`
- PDF: `{PDF_PATH.name}`
- Build script: `scripts/build_neuronbench_supervisor_slides.py`

The historical top-1/top-2/top5 bars use the manually audited whole-frontier
recovery rates in `reports/neuron_topk_manual_transfer_report.pdf`. The newest
run 708907 panel uses strict numerical recovery (`NRMSE <= 1e-6`) because that
run has not yet had the same manual equation audit. “top5” refers to the one
completed leave-one-world-out fold (train on five, hold out Z-rebound), not all
six folds.
"""
    readme.write_text(text, encoding="utf-8")


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    slides = build_slides(load_newest_summary())
    with PdfPages(PDF_PATH, metadata={"Title": "NeuronBench × evolve_pysr", "Author": "meta_sr project"}) as pdf:
        for i, slide in enumerate(slides, 1):
            render_matplotlib(slide, i, pdf)
    render_pptx(slides)
    write_readme()
    print(f"Wrote {PPTX_PATH}")
    print(f"Wrote {PDF_PATH}")
    print(f"Rendered {len(slides)} slides to {ASSET_DIR}")


if __name__ == "__main__":
    main()
