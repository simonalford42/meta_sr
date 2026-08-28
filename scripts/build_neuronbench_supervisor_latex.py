#!/usr/bin/env python3
"""Build scientific figures and rendered PPTX for the LaTeX supervisor deck.

This script reads existing project artifacts only. It does not launch PySR or
submit any SLURM work. The Beamer source itself lives in presentations/.
"""

from __future__ import annotations

import argparse
import json
import shutil
import subprocess
import urllib.request
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
PRESENTATIONS = ROOT / "presentations"
STEM = "neuronbench_evolve_pysr_supervisor_2026-08-28"
TEX = PRESENTATIONS / f"{STEM}.tex"
PDF = PRESENTATIONS / f"{STEM}.pdf"
PPTX = PRESENTATIONS / f"{STEM}.pptx"
ASSETS = PRESENTATIONS / f"{STEM}_latex_assets"
RENDERED = PRESENTATIONS / f"{STEM}_assets"

HH_FIGURE = Path("/home/sca63/sandia_spiking/Hodgkins-Huxley/hh_fig1.pdf")
MURPHY_FIGURE_URL = "https://arxiv.org/html/2608.09696v4/x3.svg"

WORLD_ORDER = ["z_rebound", "h_sag", "na_fatigue", "ca_rebound", "d_type", "textbook_M"]
WORLD_LABELS = {
    "z_rebound": "Z-rebound",
    "h_sag": "H-sag",
    "na_fatigue": "Na-fatigue",
    "ca_rebound": "Ca-rebound",
    "d_type": "D-type",
    "textbook_M": "Textbook-M",
}


def style() -> None:
    plt.rcParams.update(
        {
            "font.family": "DejaVu Sans",
            "font.size": 9,
            "axes.titlesize": 10,
            "axes.labelsize": 9,
            "axes.spines.top": False,
            "axes.spines.right": False,
            "axes.grid": True,
            "grid.alpha": 0.22,
            "legend.frameon": False,
            "savefig.bbox": "tight",
            "savefig.pad_inches": 0.04,
        }
    )


def ensure_primary_figures() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(HH_FIGURE, ASSETS / "sandia_hh_fig1.pdf")
    svg = ASSETS / "murphy_fig4a.svg"
    png = ASSETS / "murphy_fig4a.png"
    if not svg.exists():
        urllib.request.urlretrieve(MURPHY_FIGURE_URL, svg)
    # CairoSVG correctly handles the CSS/use elements in the arXiv SVG.
    import cairosvg

    cairosvg.svg2png(url=str(svg), write_to=str(png), output_width=1400)


def raw_h_sag_plot() -> None:
    source = ROOT / "runs/neuronbench_fully_observable/data/h_sag.npz"
    with np.load(source) as arrays:
        X = np.asarray(arrays["X_train"], dtype=float)
        y = np.asarray(arrays["y_train"], dtype=float)

    i_ext, voltage, phi_na, phi_k, phi_h = X.T
    terms = {
        r"$I_{\rm ext}$": i_ext,
        r"$I_{\rm Na}$": -120.0 * phi_na * (voltage - 50.0),
        r"$I_{\rm K}$": -36.0 * phi_k * (voltage + 77.0),
        r"$I_{\rm leak}$": -0.3 * (voltage + 54.4),
        r"$I_h$": -5.0 * phi_h * (voltage + 30.0),
    }

    fig, axes = plt.subplots(1, 2, figsize=(8.8, 3.05), gridspec_kw={"width_ratios": [1.18, 1]})
    scatter = axes[0].scatter(voltage, y, c=phi_h, s=11, cmap="viridis", alpha=0.75, linewidths=0)
    axes[0].set_xlabel(r"membrane voltage $V$ (mV)")
    axes[0].set_ylabel(r"target $\dot V$ (mV/ms)")
    axes[0].set_title("H-sag training states (all 1,024 points)")
    cb = fig.colorbar(scatter, ax=axes[0], pad=0.02)
    cb.set_label(r"observed $\phi_h$")

    values = [terms[name] for name in terms]
    axes[1].boxplot(
        values,
        vert=False,
        labels=list(terms),
        showfliers=False,
        whis=(5, 95),
        patch_artist=True,
        boxprops={"facecolor": "#d8e6ef", "edgecolor": "#345b75"},
        medianprops={"color": "#8f3b2f", "linewidth": 1.3},
        whiskerprops={"color": "#345b75"},
        capprops={"color": "#345b75"},
    )
    axes[1].set_xscale("symlog", linthresh=10)
    axes[1].set_xlabel("additive contribution to $\\dot V$ (mV/ms; symlog)")
    axes[1].set_title("Ground-truth current contributions (5--95%)")
    fig.tight_layout(w_pad=1.0)
    fig.savefig(ASSETS / "raw_h_sag.pdf")
    fig.savefig(ASSETS / "raw_h_sag.png", dpi=220)
    plt.close(fig)


def historical_transfer_plot() -> None:
    comparison_path = ROOT / "reports/neuron_manual_match_comparison.json"
    with comparison_path.open(encoding="utf-8") as handle:
        comparison = json.load(handle)
    records = comparison["records"]

    top5_path = ROOT / "runs/312780/neuron_results.json"
    with top5_path.open(encoding="utf-8") as handle:
        top5_payload = json.load(handle)
    top5 = [run for run in top5_payload["runs"] if run["world"] == "z_rebound"]

    panels = [
        ("top-1: train on Z-rebound", "top-1", ["h_sag", "na_fatigue", "ca_rebound", "d_type", "textbook_M"]),
        ("top-2: train on Z-rebound + H-sag", "top-2", ["na_fatigue", "ca_rebound", "d_type", "textbook_M"]),
        ('"top5": held-out Z-rebound fold', "top5", ["z_rebound"]),
    ]
    fig, axes = plt.subplots(1, 3, figsize=(10.0, 3.2), sharey=True, gridspec_kw={"width_ratios": [5, 4, 1.7]})
    rng = np.random.default_rng(260828)
    for ax, (title, regime, worlds) in zip(axes, panels):
        for xi, world in enumerate(worlds):
            if regime == "top5":
                rows = [
                    {"best_nrmse": r["best_nrmse"], "manual_match": True}
                    for r in top5
                ]
            else:
                rows = [r for r in records if r["training_regime"] == regime and r["world"] == world]
            for row in rows:
                jitter = rng.uniform(-0.12, 0.12)
                matched = bool(row["manual_match"])
                ax.scatter(
                    xi + jitter,
                    float(row["best_nrmse"]),
                    s=28,
                    facecolor="#255f85" if matched else "white",
                    edgecolor="#255f85",
                    linewidth=1.0,
                    zorder=3,
                )
        ax.axhline(1e-6, color="#a23b32", linestyle="--", linewidth=1.0)
        ax.axhline(1e-3, color="#777777", linestyle=":", linewidth=0.9)
        ax.set_yscale("log")
        ax.set_title(title)
        ax.set_xticks(range(len(worlds)), [WORLD_LABELS[w].replace("-", "\n", 1) for w in worlds])
        ax.tick_params(axis="x", labelsize=7.5)
        ax.grid(axis="x", visible=False)
    axes[0].set_ylabel("best held-out NRMSE on saved frontier")
    axes[0].text(-0.45, 1.3e-6, r"strict $10^{-6}$", color="#a23b32", fontsize=7)
    axes[0].text(-0.45, 1.3e-3, r"near-exact $10^{-3}$", color="#666666", fontsize=7)
    fig.text(0.50, 0.01, "filled = manual structural match; open = no manual match", ha="center", fontsize=8)
    fig.tight_layout(rect=(0, 0.06, 1, 1), w_pad=0.8)
    fig.savefig(ASSETS / "historical_transfer_nrmse.pdf")
    fig.savefig(ASSETS / "historical_transfer_nrmse.png", dpi=220)
    plt.close(fig)


def newest_results_plot() -> None:
    result_path = ROOT / "runs/708907/neuron_full_eval/neuron_results.json"
    with result_path.open(encoding="utf-8") as handle:
        payload = json.load(handle)
    data_path = ROOT / "runs/708907/run_data.json"
    with data_path.open(encoding="utf-8") as handle:
        run_data = json.load(handle)

    worlds = ["h_sag", "na_fatigue", "ca_rebound", "d_type", "textbook_M"]
    fig, axes = plt.subplots(1, 2, figsize=(9.2, 3.15), gridspec_kw={"width_ratios": [1.55, 1]})
    rng = np.random.default_rng(708907)
    for xi, world in enumerate(worlds):
        rows = [r for r in payload["runs"] if r["world"] == world]
        for row in rows:
            assessment = row["assessment"]
            color = "#2f7d4b" if assessment == "recovered" else "#4f79a7"
            axes[0].scatter(xi + rng.uniform(-0.11, 0.11), row["best_nrmse"], s=34, color=color, edgecolor="white", linewidth=0.6, zorder=3)
    axes[0].set_yscale("log")
    axes[0].axhline(1e-6, color="#a23b32", linestyle="--", linewidth=1.0, label=r"recovered: $10^{-6}$")
    axes[0].axhline(1e-3, color="#777777", linestyle=":", linewidth=0.9, label=r"near-exact: $10^{-3}$")
    axes[0].set_xticks(range(len(worlds)), [WORLD_LABELS[w].replace("-", "\n", 1) for w in worlds])
    axes[0].tick_params(axis="x", labelsize=7.5)
    axes[0].set_ylabel("best held-out NRMSE")
    axes[0].set_title("All 25 held-out fits (five seeds per world)")
    axes[0].legend(loc="lower left", fontsize=7)
    axes[0].grid(axis="x", visible=False)

    generations = [int(g["generation"]) for g in run_data["generations"]]
    scores = [float(g["best_score"]) for g in run_data["generations"]]
    axes[1].plot(generations, scores, marker="o", markersize=3.5, color="#255f85", linewidth=1.5)
    axes[1].axhline(0.0, color="#a23b32", linestyle="--", linewidth=0.9, label="baseline")
    axes[1].set_xlim(1, 15)
    axes[1].set_ylim(-0.04, 1.06)
    axes[1].set_xlabel("generation")
    axes[1].set_ylabel("best training solve fraction")
    axes[1].set_title("Run 708907 evolution trace")
    axes[1].set_yticks([0, 0.5, 1.0])
    axes[1].legend(loc="lower right", fontsize=7)
    fig.tight_layout(w_pad=1.2)
    fig.savefig(ASSETS / "uninformative_raw_results.pdf")
    fig.savefig(ASSETS / "uninformative_raw_results.png", dpi=220)
    plt.close(fig)


def render_pdf_to_pptx() -> None:
    RENDERED.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [
            "gs",
            "-q",
            "-dNOSAFER",
            "-sDEVICE=pngalpha",
            "-r180",
            f"-o{RENDERED / 'slide_%02d.png'}",
            str(PDF),
        ],
        check=True,
    )
    pages = sorted(RENDERED.glob("slide_*.png"))
    if len(pages) != 10:
        raise RuntimeError(f"Expected 10 rendered pages, found {len(pages)}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for page in pages:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(page), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(PPTX)
    subprocess.run(
        [
            "montage",
            *[str(p) for p in pages],
            "-thumbnail",
            "640x360",
            "-tile",
            "4x3",
            "-geometry",
            "+12+12",
            str(RENDERED / "contact_sheet.png"),
        ],
        check=True,
    )


def compile_tex() -> None:
    tectonic = shutil.which("tectonic")
    if tectonic is None:
        raise RuntimeError("tectonic not found")
    subprocess.run([tectonic, TEX.name], cwd=PRESENTATIONS, check=True)


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--figures-only", action="store_true")
    args = parser.parse_args()
    style()
    ensure_primary_figures()
    raw_h_sag_plot()
    historical_transfer_plot()
    newest_results_plot()
    if not args.figures_only:
        compile_tex()
        render_pdf_to_pptx()
    print(f"Assets: {ASSETS}")
    if not args.figures_only:
        print(f"PDF: {PDF}")
        print(f"PPTX: {PPTX}")


if __name__ == "__main__":
    main()
